"""
MemoraGraph – Text Extraction from Multiple Document Formats

Supports: PDF, DOCX, TXT, CSV, JSON, PPTX
"""

import csv
import io
import json
import logging
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".txt", ".csv", ".json", ".pptx", ".md"}


def extract_text(file_path: str) -> tuple[str, dict]:
    """
    Extract raw text and basic metadata from a document.
    
    Returns:
        (text: str, metadata: dict)
    
    Raises:
        ValueError: for unsupported file types
        RuntimeError: for extraction failures
    """
    path = Path(file_path)
    ext = path.suffix.lower()

    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported file type: {ext}. Supported: {SUPPORTED_EXTENSIONS}")

    logger.info("Extracting text from: %s (type=%s)", path.name, ext)

    try:
        if ext == ".pdf":
            return _extract_pdf(path)
        elif ext == ".docx":
            return _extract_docx(path)
        elif ext in (".txt", ".md"):
            return _extract_txt(path)
        elif ext == ".csv":
            return _extract_csv(path)
        elif ext == ".json":
            return _extract_json(path)
        elif ext == ".pptx":
            return _extract_pptx(path)
    except Exception as e:
        raise RuntimeError(f"Failed to extract text from {path.name}: {e}") from e


def _extract_pdf(path: Path) -> tuple[str, dict]:
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    pages = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            pages.append(text.strip())
    full_text = "\n\n".join(pages)
    
    meta = reader.metadata or {}
    metadata = {
        "author": str(meta.get("/Author", "")) or None,
        "title": str(meta.get("/Title", "")) or None,
        "creation_date": str(meta.get("/CreationDate", "")) or None,
        "page_count": len(reader.pages),
    }
    return _clean_text(full_text), {k: v for k, v in metadata.items() if v}


def _extract_docx(path: Path) -> tuple[str, dict]:
    from docx import Document
    doc = Document(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    
    # Also extract tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)
    
    full_text = "\n\n".join(paragraphs)
    
    props = doc.core_properties
    metadata = {
        "author": props.author or None,
        "title": props.title or None,
        "created": str(props.created) if props.created else None,
        "modified": str(props.modified) if props.modified else None,
        "subject": props.subject or None,
    }
    return _clean_text(full_text), {k: v for k, v in metadata.items() if v}


def _extract_txt(path: Path) -> tuple[str, dict]:
    import chardet
    raw = path.read_bytes()
    detected = chardet.detect(raw)
    encoding = detected.get("encoding", "utf-8") or "utf-8"
    text = raw.decode(encoding, errors="replace")
    return _clean_text(text), {}


def _extract_csv(path: Path) -> tuple[str, dict]:
    import chardet
    raw = path.read_bytes()
    detected = chardet.detect(raw)
    encoding = detected.get("encoding", "utf-8") or "utf-8"
    text = raw.decode(encoding, errors="replace")
    
    reader = csv.DictReader(io.StringIO(text))
    rows = []
    for row in reader:
        row_parts = [f"{k}: {v}" for k, v in row.items() if v]
        if row_parts:
            rows.append(", ".join(row_parts))
    
    full_text = "\n".join(rows)
    return _clean_text(full_text), {"format": "csv"}


def _extract_json(path: Path) -> tuple[str, dict]:
    with open(path, encoding="utf-8", errors="replace") as f:
        data = json.load(f)
    
    def flatten(obj, prefix="") -> list[str]:
        parts = []
        if isinstance(obj, dict):
            for k, v in obj.items():
                parts.extend(flatten(v, f"{prefix}{k}: " if prefix else f"{k}: "))
        elif isinstance(obj, list):
            for i, item in enumerate(obj):
                parts.extend(flatten(item, f"{prefix}[{i}] "))
        else:
            parts.append(f"{prefix}{obj}")
        return parts
    
    lines = flatten(data)
    full_text = "\n".join(lines)
    return _clean_text(full_text), {"format": "json"}


def _extract_pptx(path: Path) -> tuple[str, dict]:
    from pptx import Presentation
    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
    
    full_text = "\n\n".join(slides)
    core_props = prs.core_properties
    metadata = {
        "author": core_props.author or None,
        "title": core_props.title or None,
        "slide_count": len(prs.slides),
    }
    return _clean_text(full_text), {k: v for k, v in metadata.items() if v}


def _clean_text(text: str) -> str:
    """Remove excess whitespace while preserving paragraph structure."""
    import re
    # Normalize line endings
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    # Remove non-printable characters (keep newlines/tabs)
    text = re.sub(r"[^\S\n\t]+", " ", text)
    # Collapse multiple blank lines to max 2
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()
